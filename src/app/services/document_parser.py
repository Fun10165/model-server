# src/app/services/document_parser.py
import io
import logging
from typing import List, Tuple, IO, TypedDict, Generator
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed


# 导入文档解析库
import docx
from pptx import Presentation
import openpyxl
import fitz  # PyMuPDF

# 导入我们已有的模型处理服务
from . import model_interactor

class DocumentParsingError(Exception):
    pass

# --- Helper Functions ---

# VVVV  NEW HELPER FUNCTION  VVVV
def _add_separators_for_non_ascii(text: str, limit: int = 2000) -> str:
    """
    Iterates through text, inserting a separator after 'limit' consecutive non-ASCII characters.
    """
    new_parts = []
    non_ascii_counter = 0
    for char in text:
        new_parts.append(char)
        if ord(char) >= 128:  # Character is non-ASCII
            non_ascii_counter += 1
            if non_ascii_counter >= limit:
                new_parts.append("\n----\n")
                non_ascii_counter = 0
        else:  # Character is ASCII, so it breaks the sequence
            non_ascii_counter = 0
    return "".join(new_parts)
# ^^^^ END OF NEW HELPER FUNCTION ^^^^

class DocumentUnit(TypedDict):
    unit_identifier: str
    text: str
    images: List[bytes]

def _chunk_list(data: list, size: int) -> Generator[list, None, None]:
    """一个辅助函数，将列表分割成指定大小的块"""
    for i in range(0, len(data), size):
        yield data[i:i + size]

# --- 提取函数部分（无改动） ---

def _extract_from_pptx(file_stream: IO[bytes]) -> List[DocumentUnit]:
    presentation = Presentation(file_stream)
    units: List[DocumentUnit] = []
    for i, slide in enumerate(presentation.slides):
        slide_text = "\n".join([shape.text for shape in slide.shapes if shape.has_text_frame]).strip()
        images = [shape.image.blob for shape in slide.shapes if hasattr(shape, "image")]
        units.append({
            "unit_identifier": f"幻灯片 {i + 1}",
            "text": slide_text,
            "images": images
        })
    logging.info(f"从PPTX中提取了 {len(units)} 个幻灯片单元。")
    return units

def _extract_from_pdf(file_stream: IO[bytes]) -> List[DocumentUnit]:
    doc = fitz.open(stream=file_stream, filetype="pdf")
    units: List[DocumentUnit] = []
    for i, page in enumerate(doc):
        page_text = page.get_text("text").strip()
        images_on_page = []
        for img in page.get_images(full=True):
            xref = img[0]
            base_image = doc.extract_image(xref)
            images_on_page.append(base_image["image"])
        units.append({
            "unit_identifier": f"页面 {i + 1}",
            "text": page_text,
            "images": images_on_page
        })
    logging.info(f"从PDF中提取了 {len(units)} 个页面单元。")
    return units

def _extract_from_docx(file_stream: IO[bytes]) -> List[DocumentUnit]:
    document = docx.Document(file_stream)
    full_text = "\n".join([para.text for para in document.paragraphs]).strip()
    images = [rel.target_part.blob for rel in document.part.rels.values() if "image" in rel.target_ref]
    unit: DocumentUnit = { "unit_identifier": "文档内容", "text": full_text, "images": images }
    logging.info(f"从DOCX中提取了 1 个文档单元，包含 {len(images)} 张图片。")
    return [unit]

def _extract_from_xlsx(file_stream: IO[bytes]) -> List[DocumentUnit]:
    workbook = openpyxl.load_workbook(file_stream)
    units: List[DocumentUnit] = []
    for sheet in workbook.worksheets:
        unit_identifier = f"工作表 '{sheet.title}'"
        sheet_text = "\n".join(str(cell.value) for row in sheet.iter_rows() for cell in row if cell.value is not None).strip()
        images = [image.ref for image in sheet._images] if hasattr(sheet, '_images') else []
        if sheet_text or images:
             units.append({ "unit_identifier": unit_identifier, "text": sheet_text, "images": images })
    logging.info(f"从XLSX中提取了 {len(units)} 个工作表单元。")
    return units

def process_document_images(file_content: bytes, filename: str) -> str:
    # ... (The entire main logic from the previous step remains exactly the same) ...
    # ... (This includes the dual-branch logic for "direct output" vs "aggregate" modes) ...
    # ...
    # After the `if use_direct_output_mode: ... else: ...` block, you will have the `final_content_parts` list.
    # The code below should be placed right after that block.
    # ...
    file_ext = filename.split('.')[-1].lower()
    file_stream = io.BytesIO(file_content)

    extraction_map = {
        'pptx': _extract_from_pptx, 'pdf': _extract_from_pdf,
        'docx': _extract_from_docx, 'xlsx': _extract_from_xlsx,
    }
    if file_ext not in extraction_map:
        raise DocumentParsingError(f"不支持的文件类型。当前支持 'pptx', 'pdf', 'docx', 'xlsx'。")

    document_units = extraction_map[file_ext](file_stream)
    
    if not any(unit['images'] for unit in document_units):
        logging.info("文档中未找到图片，将返回提取的文本内容。")
        all_texts = [f"--- {unit['unit_identifier']} ---\n{unit['text']}" for unit in document_units if unit['text']]
        return f"文档中未找到图片，返回提取的文本内容。\n\n" + "\n\n".join(all_texts)

    units_with_images = [u for u in document_units if u['images']]
    is_unpaginated = file_ext in ('docx', 'xlsx')
    
    use_direct_output_mode = is_unpaginated or len(units_with_images) <= 10
    
    MAX_WORKERS = 5
    final_content_parts = []
    
    if use_direct_output_mode:
        logging.info(f"激活 [直接输出模式]，因文档类型为 '{file_ext}' 或图片单元数 ({len(units_with_images)}) 不超过10。")
        DIRECT_ANALYSIS_PROMPT = "你是一个专业的图像分析助手。请详细描述我提供给你的这一批图片的内容。你的描述应该客观、详尽，并使用Markdown格式。不要做任何与图片无关的推断。"
        
        analysis_tasks = []
        for unit in units_with_images:
            for i, image_chunk in enumerate(_chunk_list(unit['images'], 8)):
                task_info = (unit['unit_identifier'], i)
                analysis_tasks.append((DIRECT_ANALYSIS_PROMPT, image_chunk, task_info))

        per_unit_analysis_results = defaultdict(list)
        with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            future_to_info = {
                executor.submit(model_interactor.get_model_response, prompt, img_bytes, model="doubao-seed-1-6-flash-250715"): info
                for prompt, img_bytes, info in analysis_tasks
            }
            for future in as_completed(future_to_info):
                unit_identifier, chunk_index = future_to_info[future]
                try:
                    result = future.result()
                    per_unit_analysis_results[unit_identifier].append((chunk_index, result))
                except Exception as e:
                    error_message = f"批次 {chunk_index} 分析失败: {e}"
                    per_unit_analysis_results[unit_identifier].append((chunk_index, error_message))

        for unit in document_units:
            identifier = unit['unit_identifier']
            text_part = unit.get('text', '').strip()
            final_content_parts.append(f"### {identifier}\n\n{text_part if text_part else '此单元无文本内容。'}")
            if unit['images']:
                analyses = per_unit_analysis_results.get(identifier)
                if analyses:
                    analyses.sort(key=lambda x: x[0])
                    analysis_texts = [res for _, res in analyses]
                    combined_analysis = "\n\n".join(analysis_texts)
                    final_content_parts.append(f"\n--- 图片分析 ---\n\n{combined_analysis}")
        
    else:
        logging.info(f"激活 [聚合分析模式]，因文档类型为 '{file_ext}' 且图片单元数 ({len(units_with_images)}) 超过10。")
        STAGE1_PROMPT_TEMPLATE = "..." # As before
        analysis_tasks = []
        # ... (rest of the aggregate mode logic as before)
        STAGE1_PROMPT_TEMPLATE = """你是一个专业的文档分析助手。任务是结合我提供的“页面文字”和一组“页面上的图片”，生成对这个页面**当前批次内容**的综合性描述。你的分析应专注于当前提供的信息，并使用Markdown格式排版。尽可能保留原有的所有文字。如果你觉得图片与文字无关，那么请直接输出文字和你对图片的描述。

---
**{unit_identifier}的文字**:
{text}
---
**{unit_identifier}的图片** (本批次共{image_count}张):
[图片内容已提供，请开始分析]"""

        analysis_tasks = []
        for unit in units_with_images:
            for i, image_chunk in enumerate(_chunk_list(unit['images'], 8)):
                prompt = STAGE1_PROMPT_TEMPLATE.format(unit_identifier=unit['unit_identifier'], text=unit['text'], image_count=len(image_chunk))
                task_info = (unit['unit_identifier'], i)
                analysis_tasks.append((prompt, image_chunk, task_info))

        per_unit_analysis_results = defaultdict(list)
        with ThreadPoolExecutor(max_workers=MAX_WORKERS) as executor:
            future_to_info = {
                executor.submit(model_interactor.get_model_response, prompt, img_bytes, model="doubao-seed-1-6-flash-250715"): info
                for prompt, img_bytes, info in analysis_tasks
            }
            for future in as_completed(future_to_info):
                unit_identifier, _ = future_to_info[future]
                try:
                    per_unit_analysis_results[unit_identifier].append(future.result())
                except Exception as e:
                    per_unit_analysis_results[unit_identifier].append(f"分析失败: {e}")
        STAGE2_AGGREGATION_PROMPT = """你是一个高级内容编辑。你收到了针对【同一个文档页面/单元】的几份独立分析报告，这是因为该页面的图片过多被分批处理了。你的任务是将这些分散的、描述同一页面的报告，合并成一份最终的、连贯的、完整的分析。

请严格遵循以下规则：
1.  **无缝合并**：消除报告间的重复引言和割裂感，将内容融合成一段流畅的文本。
2.  **保留所有细节**：确保原始报告中提到的所有图片分析和文本要点都被包含在最终版本中。尽可能保留原有的所有文字。
3.  **单一视角**：最终的输出应该读起来像是对这一个页面的一次性完整分析。
4.  **Markdown格式**: 使用合适的Markdown标题和格式来组织最终报告。

---
**来自【{unit_identifier}】的多份独立分析报告**:
{combined_analyses}
---
请根据以上报告，为【{unit_identifier}】生成最终的、合并后的一份综合性分析。"""

        for unit in document_units:
            identifier = unit['unit_identifier']
            if not unit['images']:
                if unit['text']: final_content_parts.append(f"--- {identifier} (无图片) ---\n{unit['text']}")
                continue
            
            analyses = per_unit_analysis_results.get(identifier, [])
            if len(analyses) == 1:
                final_content_parts.append(f"### 对 {identifier} 的分析\n{analyses[0]}")
            elif len(analyses) > 1:
                combined_analyses = "\n\n---\n".join(analyses)
                final_prompt = STAGE2_AGGREGATION_PROMPT.format(unit_identifier=identifier, combined_analyses=combined_analyses)
                try:
                    aggregated_result = model_interactor.get_model_response(prompt=final_prompt, model="doubao-seed-1-6-flash-250715")
                    final_content_parts.append(aggregated_result)
                except Exception as e:
                    logging.error(f"单元 '{identifier}' 的聚合步骤失败: {e}")
                    error_header = f"### 对 {identifier} 的分析 (聚合失败)\n\n**错误**: 未能将以下分批报告合并成最终版本。已为您呈现原始报告：\n\n---"
                    final_content_parts.append(error_header + "\n\n" + combined_analyses)

    # --- 最终返回前的后处理步骤 ---
    # VVVV  NEW POST-PROCESSING LOGIC  VVVV
    
    # 1. 组合成初步的最终报告
    initial_report = "\n\n".join(final_content_parts).strip()
    
    # 2. 为长串的非ASCII字符添加分隔符
    logging.info("正在检查并为长串非ASCII字符添加分隔符...")
    report_with_separators = _add_separators_for_non_ascii(initial_report)
    
    # 3. 检查最终长度并按需进行总结
    final_length = len(report_with_separators)
    logging.info(f"报告在总结前的最终长度为: {final_length} 字符。")
    
    if final_length > 8000:
        logging.warning(f"报告长度 ({final_length}) 超出8000字符限制，将启动最终总结...")
        
        SUMMARIZATION_PROMPT = """你是一个专业的文档摘要助手。以下是一份详细的文档分析报告，但它对于最终系统来说太长了。
你的任务是：
1.  **精确总结**：将报告内容进行浓缩，同时必须保留所有核心的发现、关键数据、重要结论和图片分析的要点。
2.  **控制长度**：最终的总结文本**绝对不能**超过8000个字符。
3.  **保持格式**：尽可能地保留原始报告中的Markdown格式（如标题、列表）以保证可读性。

请开始对下面的报告进行总结：
---
"""
        try:
            summarized_report = model_interactor.get_model_response(
                prompt=SUMMARIZATION_PROMPT + report_with_separators,
                model="doubao-seed-1-6-flash-250715" # Or a model known for good summarization
            )
            logging.info(f"总结完成。新长度: {len(summarized_report)} 字符。")
            return summarized_report
        except Exception as e:
            logging.error(f"最终总结步骤失败: {e}。将返回被截断的原始报告。")
            # 降级方案：返回一个被截断的版本，以确保不会超出硬性限制
            return report_with_separators[:7990] + "\n\n[...报告过长且总结失败，已被截断...]"
            
    else:
        logging.info("报告长度在限制范围内，无需总结。")
        return report_with_separators
    # ^^^^ END OF NEW POST-PROCESSING LOGIC ^^^^